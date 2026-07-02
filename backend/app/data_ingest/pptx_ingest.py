# data_ingest/pptx_ingest.py
# Parses PowerPoint slides (.pptx) using the python-pptx library.

import logging
from typing import List

logger = logging.getLogger(__name__)

def load_documents_from_pptx(file_path: str) -> List[str]:
    """
    Extracts text shapes slide-by-slide from a PowerPoint presentation.
    """
    slide_texts = []
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_content = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                    
            if slide_content:
                # Merge slide shape texts with a slide identifier
                slide_texts.append(f"Slide {slide_num}: " + " | ".join(slide_content))
                
        logger.info(f"Extracted {len(slide_texts)} slides from PPTX: {file_path}")
    except Exception as exc:
        logger.error(f"Error parsing PPTX file {file_path}: {exc}")
        raise
        
    return slide_texts
